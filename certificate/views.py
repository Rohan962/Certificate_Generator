from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import Event, Participant
import pandas as pd
from .convter import ppt2pdf
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os
import sys
import re
import time
import subprocess
import boto3

# from comtypes.client import CreateObject
# import comtypes
# from comtypes import CoInitialize, CoUninitialize

# Initialize the SES client
ses_client = boto3.client('ses', region_name='us-east-1')  # Replace 'your-region' with your AWS SES region

def index(request):
	return render(request, 'index.html')

@login_required
def create(request):
	if request.method == "POST":
		csv = request.FILES.get('csv')
		temp = request.FILES.get('template')
		# print(f"template file :  {temp}")
		event_name = request.POST.get('event_name')
		
		event = Event(user = request.user,
			event_name = event_name,
			csv_file = csv,
			template = temp)
		print(f"template file :  {event.template}")
		event.save()

		return redirect(f"/certificate/{event.id}/{event.slug}")

	return render(request, 'certificate/create_event.html')

@login_required
def delete_event(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
 
	
	if event.user == request.user:
	    event.delete()
	return redirect('view_certificate_status')

@login_required
def track(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	# CoInitialize()
	if event.message:

		return render(request, 'certificate/track.html', {
			'event_name': event.event_name,
			'event_date': event.date,
			'participat_details': Participant.objects.filter(event=event)
			})
	
	def convert_ppt_to_pdf(ppt_path, pdf_path):
			# powerpoint = CreateObject("PowerPoint.Application")
			# powerpoint.Visible = 1
			# try:
			# 	presentation = powerpoint.Presentations.Open(ppt_path)
			# 	presentation.SaveAs(pdf_path, 32)  # 32 is the format for PDF
			# 	presentation.Close()
			# except Exception as e:
			# 	print(f"Failed to convert {ppt_path} to PDF: {e}")
			# finally:
			# 	powerpoint.Quit()
		try:
			command = [
				"libreoffice", "--headless", "--convert-to", "pdf", "--outdir",
				os.path.dirname(pdf_path), ppt_path
			]
			subprocess.run(command, check=True)
			print(f"✅ Converted {ppt_path} to {pdf_path}")
		except subprocess.CalledProcessError as e:
			print(f"❌ Failed to convert {ppt_path} to PDF: {e}")
    
	def delete_generated_files(files):
		#"""Deletes files after processing"""
		time.sleep(5)  # Give some time for files to release if in use
		for file_path in files:
			if os.path.exists(file_path):
				try:
					os.remove(file_path)
					print(f"✅ Deleted: {file_path}")
				except Exception as e:
					print(f"❌ Error deleting {file_path}: {e}")
			else:
				print(f"⚠️ File not found: {file_path}")


	prs = Presentation(event.template)
	print(event.template)
	# file_prefix = re.search(r'([^\\]+)\.pptx$', event.template)
	file_prefix = os.path.splitext(os.path.basename(event.template.name))[0] 
	 

	print(file_prefix)
	st=""
	for slide in prs .slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				st = st + shape.text
				st = st + " "
	li = st.split()
	tags = []
	for i in li:
		if i[0] == "<" and i[-1] == ">":
			tags.append(i)
	
	if request.method == "POST":
		email_col = request.POST.get('emails')
		subject = request.POST.get('subject')
		mess = request.POST.get('mess')
		values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
		
		event.email_column = email_col
		event.message = mess
		event.subject = subject
		event.save()

		df=pd.read_csv(event.csv_file)
		df_len=df.shape

		# List to track generated files
		files_to_delete = []
		i=0

		data = {
        	"client_id":"1054669543742-tda6tkgc4hukhgfbomq93ajvi3iumiuq.apps.googleusercontent.com",
        	"client_secret":"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",
        	"refresh_token": "1//04UA4eT3XoKjNCgYIARAAGAQSNwF-L9Ir0NO4JvnUO_G9r3clMrFWGsw_DupMVrSiPVwu2N37X1RrO7SDEDrY9-xckg2juTev2e4",
        	'grant_type': 'refresh_token'
            }
		a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
		token = f"Bearer {dict(a.json()).get('access_token')}"
		li=["First","Second","Third"]
		while i < df_len[0]:
			prs = Presentation(event.template)
			j=""
			if i<9:
				j="00"
			elif i>=9 and i < 99 :
				j="0"
			
			for tag, v_type, value in values:
				for slide in prs.slides:
					for shape in slide.shapes:
						if shape.has_text_frame:
							text_frame = shape.text_frame
							for paragraph in text_frame.paragraphs:
								# Get full paragraph text
								full_text = "".join(run.text for run in paragraph.runs)
								print(f"Original text: {full_text}")
								# Apply all replacements in one step
								for tag, v_type, value in values:
									if tag in full_text:  # Check if tag exists before replacing
										if v_type == 'text':
											full_text = full_text.replace(tag, value)
										elif v_type == 'date':
											date_parts = value.split('/')
											if len(date_parts) == 3:  # Ensure correct format
												formatted_date = f"{date_parts[0]}/{date_parts[1]}/{date_parts[2]}"
												full_text = full_text.replace(tag, formatted_date)
											else:
												print(f"Warning: Unexpected date format for '{value}'")
										elif v_type == 'csv':
											full_text = full_text.replace(tag, df.loc[i, value])
										elif v_type == "auto":
											full_text = full_text.replace(tag, value + "/" + j + str(i + 1))
								print(f"Modified text: {full_text.replace(tag, value)}")
								# Remove old runs
								for run in paragraph.runs:
									run.text = ""  

								# Insert new text in the first run
								paragraph.runs[0].text = full_text

			base_folder = os.getcwd()  # Example base folder
			sub_folder = "media\certificates\generated"
			output_folder_path = os.path.join(base_folder, sub_folder)
				
			s_name = df.loc[i,event.email_column].split('@')[0]
			stud_name = df.loc[i,"Name"]
			 
			ppt_path = (f"{output_folder_path}\{file_prefix}_{stud_name}.pptx")
			pdf_path = (f"{output_folder_path}\{file_prefix}_{stud_name}.pdf")
			
			ppt_path = ppt_path.replace("template","")
			pdf_path = pdf_path.replace("template","")


			# Store file paths
			files_to_delete.append(ppt_path)
			files_to_delete.append(pdf_path)

			prs.save(ppt_path)
			# pdf_path = (s_name+".pdf")
   
			# print(output_folder_path)
			# sys.exit()
			convert_ppt_to_pdf(ppt_path, pdf_path)
			# f_id = ppt2pdf(s_name+".pptx",s_name, token)
			# #r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
			# r = requests.get(f"https://www.googleapis.com/drive/v3/files/{f_id}/export?mimeType=application/pdf", headers={'Authorization': token}, allow_redirects=True)

			# open(s_name+'.pdf', 'wb').write(r.content)

			# Email content with line breaks and proper HTML formatting
			EMAIL_SUBJECT = "Embrizon | Congratulations on Completing Your Industrial Training & Internship at Embrizon Technologies!"
			EMAIL_BODY_TEMPLATE = """
			Dear {name},<br><br>

			<b>Congratulations on successfully completing your internship and training at Embrizon Technologies!</b><br><br>
			Your dedication, enthusiasm, and hard work have made a lasting impact during your time with us.<br><br>
			We are thrilled to witness your growth and the valuable contributions you've made to specific projects or tasks.<br>
			Your positive attitude has set a great example for your peers, and we believe the skills you've gained will serve 
			as a strong foundation for your future endeavors.<br><br>

			As you celebrate this achievement, we encourage you to showcase your accomplishments on LinkedIn by posting 
			your certificates. Please remember to <b>Tag Embrizon Technologies</b> in your post, as this is a mandatory step 
			to share the news with our network and acknowledge your outstanding achievements.<br><br>

			We would also like to invite you to join our exclusive Telegram Community, where you can find internship and career opportunities. 
			It's a great platform for you to connect with peers, mentors, and industry professionals. Please join the Community using the following link: 
			<a href="https://t.me/+mokhPkUMbGoyMDQ1" target="_blank">Embrizon Discord Server</a><br><br>

			Thank you for your valuable contributions, and we wish you continued success in your journey. Feel free to reach out 
			if you ever need a reference or support.<br><br>

			--<br>
			<b>Thanks & Regards</b><br>
			Team Embrizon<br>
			+91-9875637467<br>
			<a href="https://www.embrizon.com">www.embrizon.com</a><br>
			"""

			try:
				# Format email body with the recipient's name
				email_body = EMAIL_BODY_TEMPLATE.format(name=stud_name)
       
				mail = EmailMessage(
        			EMAIL_SUBJECT,
					email_body,
					settings.EMAIL_HOST_USER,
					[df.loc[i,event.email_column]])
				 # Set email content as HTML
				mail.content_subtype = "html"
				mail.attach_file(pdf_path)
				mail.send()
				print("Mail sent successfully")
				Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
				
			except:
				print("Mail not sent")
				Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
				# os.remove(s_name+'.pdf')
				# os.remove(s_name+".pptx")
			i=i+1

		messages.success(request, "Certificates Sent Successfuly !!")
		# Step 4: Delete files after processing
		delete_generated_files(files_to_delete)
		# print("Press Enter to continue...")
		return redirect(f"/certificate/{event.id}/{event.slug}")

	# CoUninitialize()
 
	return render(request, 'certificate/map_tags_of_template.html',{
     
		'columns': list(pd.read_csv(event.csv_file).columns),
		'tags': tags,
		})


@login_required
def view_certificate_status(request):
	return render(request, 'certificate/view_certificate_status.html',{
		'events': Event.objects.filter(user=request.user)
		})

